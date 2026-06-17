"""③ PPT 渲染智能体（见 work-person-c.md 4.3.3）。

将章节内容编排为幻灯片，使用 python-pptx 真实生成 .pptx 文件。
渲染策略：封面 → 每章一个 section slide + 内容 slide（按 markdown 列表拆分 bullets）。
"""
from __future__ import annotations

import logging
import re
from pathlib import Path
from typing import List

from ..models.dto import SectionContent, SlideData

logger = logging.getLogger(__name__)

# 幻灯片布局映射（python-pptx layout index）
SLIDE_LAYOUTS = {
    "cover": 0,       # 封面（标题幻灯片）
    "section": 1,     # 章节标题
    "content": 5,     # 标题 + 正文（使用 blank-ish 布局，手动加文本框）
    "comparison": 5,
    "code": 5,
    "summary": 5,
}


class PptRenderer:
    """PPT 渲染智能体。"""

    def arrange_slides(self, title: str, sections: List[SectionContent]) -> List[SlideData]:
        """将章节内容编排为幻灯片数据序列。"""
        slides: List[SlideData] = [SlideData(layout="cover", title=title, bullets=[])]
        for section in sections:
            slides.append(
                SlideData(layout="section", title=section.title, bullets=[], notes="章节封面")
            )
            bullets = self._markdown_to_bullets(section.markdown)
            # 每张内容 slide 最多 6 条要点
            for i in range(0, max(len(bullets), 1), 6):
                chunk = bullets[i : i + 6] or ["（本节内容）"]
                slides.append(
                    SlideData(
                        layout="content",
                        title=section.title,
                        bullets=chunk,
                    )
                )
        slides.append(SlideData(layout="summary", title="总结", bullets=["感谢学习！", "欢迎复习本套资料。"]))
        return slides

    def render(self, title: str, slides: List[SlideData], output_dir: Path) -> Path:
        """渲染为 .pptx 文件，返回文件路径。"""
        from pptx import Presentation
        from pptx.util import Inches, Pt

        prs = Presentation()
        # 统一使用宽屏比例
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        for slide_data in slides:
            layout_idx = SLIDE_LAYOUTS.get(slide_data.layout, 5)
            layout = prs.slide_layouts[layout_idx]
            slide = prs.slides.add_slide(layout)

            # 标题（部分布局自带标题占位符）
            if slide.shapes.title is not None:
                slide.shapes.title.text = slide_data.title
            else:
                box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.5), Inches(1.2))
                box.text_frame.text = slide_data.title
                for para in box.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.size = Pt(36)
                        run.font.bold = True

            # 正文要点
            if slide_data.bullets:
                body = slide.shapes.add_textbox(
                    Inches(0.8), Inches(2.0), Inches(11.5), Inches(5.0)
                )
                tf = body.text_frame
                tf.word_wrap = True
                for i, bullet in enumerate(slide_data.bullets):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = f"• {bullet}"
                    for run in p.runs:
                        run.font.size = Pt(22)

        safe_title = re.sub(r"[\\/:*?\"<>|\x00-\x1f]+", "_", title).strip("_ ")
        file_path = output_dir / f"{safe_title}.pptx"
        prs.save(str(file_path))
        logger.info("PPT 已渲染：%s（%d 页）", file_path, len(slides))
        return file_path

    @staticmethod
    def _markdown_to_bullets(markdown: str) -> List[str]:
        """将 markdown 提取为要点列表（取 `-`/`*` 列表项与小标题行）。"""
        bullets: List[str] = []
        for line in markdown.splitlines():
            stripped = line.strip()
            if not stripped:
                continue
            if stripped.startswith(("## ", "### ")):
                # 小标题转为带强调的要点
                bullets.append(stripped.lstrip("# ").strip())
            elif re.match(r"^[-*]\s+", stripped):
                bullets.append(re.sub(r"^[-*]\s+", "", stripped))
            elif len(stripped) < 80 and not stripped.startswith((">", "```")):
                bullets.append(stripped)
        return bullets[:18]  # 控制总量
